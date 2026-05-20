"""
Shared annotation primitives for reviewer-workflow figures.

The chapter's figures use one style: USACE-blue highlight boxes around the
elements being called out, with numbered circles at one corner of each box.
The text legend mapping numbers to descriptions lives in the figure caption
(via the <Figure> component's caption prop), not on the image itself —
keeping image annotation uncluttered no matter how dense the page is.
"""

from pathlib import Path
from PIL import Image, ImageDraw, ImageFont

ACCENT = (74, 124, 155, 255)        # USACE blue
ACCENT_DARK = (50, 90, 115, 255)
WHITE = (255, 255, 255, 255)
BOX_W = 4
CIRCLE_R = 18


def load_font(size: int) -> ImageFont.ImageFont:
    for path in ("C:/Windows/Fonts/segoeuib.ttf", "C:/Windows/Fonts/arialbd.ttf"):
        if Path(path).exists():
            return ImageFont.truetype(path, size)
    return ImageFont.load_default()


def draw_callout(draw: ImageDraw.ImageDraw, x: int, y: int, w: int, h: int, corner: str, idx: int, font_circle: ImageFont.ImageFont) -> None:
    """
    Draw a highlight box around (x, y, w, h) and place a numbered circle at
    the chosen corner ("tl", "tr", "bl", or "br") so it half-overlaps the
    box edge.
    """
    draw.rectangle([x - 2, y - 2, x + w + 2, y + h + 2], outline=ACCENT, width=BOX_W)

    if corner == "tl":
        cx, cy = x - 2, y - 2
    elif corner == "tr":
        cx, cy = x + w + 2, y - 2
    elif corner == "bl":
        cx, cy = x - 2, y + h + 2
    else:
        cx, cy = x + w + 2, y + h + 2

    draw.ellipse(
        [cx - CIRCLE_R, cy - CIRCLE_R, cx + CIRCLE_R, cy + CIRCLE_R],
        fill=ACCENT, outline=WHITE, width=3,
    )
    text = str(idx)
    bb = draw.textbbox((0, 0), text, font=font_circle)
    tw, th = bb[2] - bb[0], bb[3] - bb[1]
    draw.text((cx - tw // 2 - bb[0], cy - th // 2 - bb[1]), text, font=font_circle, fill=WHITE)


def mask_region(draw: ImageDraw.ImageDraw, x: int, y: int, w: int, h: int, color: tuple = WHITE) -> None:
    """
    Paint over an arbitrary region. Used to hide GitHub UI detritus that
    leaked past the capture script's hygiene CSS (announcement toasts,
    "what's new" banners, etc.).
    """
    draw.rectangle([x, y, x + w, y + h], fill=color)


def annotate_and_crop(src: Path, callouts: list, crop: tuple, out: Path, masks: list = None, pptx_out: Path = None) -> None:
    """
    callouts: list of (x, y, w, h, corner, idx) — coordinates relative to
              the FULL source image, before cropping.
    crop:     (left, top, right, bottom) bounding box on the source
    masks:    optional list of (x, y, w, h) regions to white-out before annotation
    pptx_out: optional Path; if provided, also write an editable .pptx
              alongside the PNG, with each highlight box and numbered
              circle as a free-floating PowerPoint shape the user can
              drag and resize.
    """
    img = Image.open(src).convert("RGBA")
    overlay = Image.new("RGBA", img.size, (0, 0, 0, 0))
    draw_mask = ImageDraw.Draw(img)
    draw_over = ImageDraw.Draw(overlay)
    font_circle = load_font(20)

    for region in masks or []:
        mask_region(draw_mask, *region)

    for (x, y, w, h, corner, idx) in callouts:
        draw_callout(draw_over, x, y, w, h, corner, idx, font_circle)

    composed = Image.alpha_composite(img, overlay).convert("RGB")
    cropped = composed.crop(crop)
    out.parent.mkdir(parents=True, exist_ok=True)
    cropped.save(out, "PNG", optimize=True)
    print(f"Wrote {out} ({cropped.size[0]}x{cropped.size[1]})")

    if pptx_out is not None:
        # For the .pptx, use the un-annotated cropped image as the slide
        # background so the user can freely reposition shapes without
        # baked-in marks. Shapes are layered on top.
        clean = Image.open(src).convert("RGB")
        for region in masks or []:
            ImageDraw.Draw(clean).rectangle([region[0], region[1], region[0] + region[2], region[1] + region[3]], fill=(255, 255, 255))
        clean_cropped = clean.crop(crop)
        make_pptx_overlay(clean_cropped, callouts, crop, pptx_out)


def make_pptx_overlay(cropped_img: "Image.Image", callouts: list, crop: tuple, out: Path) -> None:
    """
    Write a PowerPoint file with the cropped image as the slide background
    and one editable rounded-rectangle + numbered-oval per callout. Coords
    map from the source image's pixel space into PowerPoint EMUs by
    subtracting the crop offset and scaling at 96 DPI.

    The user can open the .pptx in PowerPoint, drag any shape or number to
    a new position, resize it, change colors, or delete it.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    import io

    crop_left, crop_top, crop_right, crop_bottom = crop
    img_w_px = crop_right - crop_left
    img_h_px = crop_bottom - crop_top
    DPI = 96
    # Slide size matches image at 96 DPI (so 1 image pixel ≈ 1/96 inch).
    slide_w = Inches(img_w_px / DPI)
    slide_h = Inches(img_h_px / DPI)

    def px_to_emu(px):
        return Emu(int(px / DPI * 914400))

    prs = Presentation()
    prs.slide_width = slide_w
    prs.slide_height = slide_h
    blank = prs.slide_layouts[6]  # fully blank layout
    slide = prs.slides.add_slide(blank)

    # Background image
    buf = io.BytesIO()
    cropped_img.save(buf, format="PNG")
    buf.seek(0)
    slide.shapes.add_picture(buf, 0, 0, slide_w, slide_h)

    accent = RGBColor(74, 124, 155)        # USACE blue
    accent_dark = RGBColor(50, 90, 115)
    white = RGBColor(255, 255, 255)

    CIRCLE_DIAM_PX = 36

    for (x, y, w, h, corner, idx) in callouts:
        # Translate from source-image space to cropped-image space.
        cx = x - crop_left
        cy = y - crop_top

        # Highlight rectangle
        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            px_to_emu(cx - 2),
            px_to_emu(cy - 2),
            px_to_emu(w + 4),
            px_to_emu(h + 4),
        )
        rect.name = f"highlight-{idx}"
        rect.fill.background()  # no fill
        rect.line.color.rgb = accent
        rect.line.width = Pt(3)

        # Numbered circle at the chosen corner
        if corner == "tl":
            ox, oy = cx - 2, cy - 2
        elif corner == "tr":
            ox, oy = cx + w + 2, cy - 2
        elif corner == "bl":
            ox, oy = cx - 2, cy + h + 2
        else:
            ox, oy = cx + w + 2, cy + h + 2

        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            px_to_emu(ox - CIRCLE_DIAM_PX / 2),
            px_to_emu(oy - CIRCLE_DIAM_PX / 2),
            px_to_emu(CIRCLE_DIAM_PX),
            px_to_emu(CIRCLE_DIAM_PX),
        )
        circle.name = f"number-{idx}"
        circle.fill.solid()
        circle.fill.fore_color.rgb = accent
        circle.line.color.rgb = white
        circle.line.width = Pt(2.5)

        tf = circle.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(idx)
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = white

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"Wrote {out}")
