"""
Storyboard GIF: filtering the Files changed view to a single commit
(Section 3 of the reviewer-workflow chapter).

Walks the reviewer through the commit-range selector:
  1. Spot the commit selector pill on the toolbar
  2. Click it — a dropdown opens
  3. Pick the content commit from the list
  4. The diff narrows to just that commit's files

The base image is the already-filtered Files changed capture. The
"open dropdown" state in frames 2-3 is a Pillow-drawn mock layered over
the page — the storyboard exists to guide the user re-recording in
ScreenToGif, not to be a true capture of every state.
"""

import io
from pathlib import Path
from PIL import Image, ImageDraw, ImageFont

REPO_ROOT = Path(__file__).resolve().parents[2]
SRC = REPO_ROOT / "planning" / "assets" / "captures" / "pr-files-changed.png"
OUT_GIF = REPO_ROOT / "static" / "figures" / "dev" / "reviewer-workflow" / "gif-commit-filter-storyboard.gif"
OUT_PPTX = REPO_ROOT / "planning" / "assets" / "figure-pptx" / "gif-commit-filter-storyboard.pptx"

CROP = (20, 200, 1340, 700)  # Tabs + toolbar + file tree + a few diff rows

ACCENT = (74, 124, 155, 255)
WHITE = (255, 255, 255, 255)
RED = (200, 30, 30, 255)
LABEL_BG = (255, 248, 220, 250)
BLUE_HOVER = (221, 244, 255, 255)


def load_font(size):
    for p in ("C:/Windows/Fonts/segoeuib.ttf", "C:/Windows/Fonts/arialbd.ttf", "C:/Windows/Fonts/segoeui.ttf"):
        if Path(p).exists():
            return ImageFont.truetype(p, size)
    return ImageFont.load_default()


def caption(draw, x, y, text, font):
    bb = draw.textbbox((0, 0), text, font=font)
    tw, th = bb[2] - bb[0], bb[3] - bb[1]
    pad = 12
    rect = [x, y, x + tw + pad * 2, y + th + pad * 2]
    draw.rounded_rectangle([rect[0] - 2, rect[1] - 2, rect[2] + 2, rect[3] + 2], radius=12, fill=WHITE)
    draw.rounded_rectangle(rect, radius=10, fill=LABEL_BG, outline=(50, 90, 115, 255), width=2)
    draw.text((x + pad - bb[0], y + pad - bb[1]), text, font=font, fill=(40, 40, 40, 255))


def draw_dropdown(draw, x, y, items, highlighted=None):
    """Draw a faux dropdown menu opening below the commit selector."""
    item_h = 44
    pad_x = 14
    w = 560
    h = item_h * len(items) + 16
    # Card
    draw.rounded_rectangle([x, y, x + w, y + h], radius=8, fill=(255, 255, 255, 255), outline=(208, 215, 222, 255), width=1)
    # Drop shadow (faint)
    f_title = load_font(13)
    f_sub = load_font(11)
    for i, (title, subtitle) in enumerate(items):
        iy = y + 8 + i * item_h
        if i == highlighted:
            draw.rectangle([x + 4, iy, x + w - 4, iy + item_h], fill=BLUE_HOVER)
        draw.text((x + pad_x, iy + 4), title, font=f_title, fill=(36, 41, 47, 255))
        draw.text((x + pad_x, iy + 22), subtitle, font=f_sub, fill=(101, 109, 118, 255))


def render_frame(base, n):
    img = base.copy()
    overlay = Image.new("RGBA", img.size, (0, 0, 0, 0))
    d = ImageDraw.Draw(overlay)
    f_cap = load_font(18)
    f_step = load_font(28)

    # Step indicator
    d.rounded_rectangle([20, 18, 130, 56], radius=8, fill=ACCENT)
    d.text((35, 24), f"Step {n}", font=f_step, fill=WHITE)

    # Commit selector pill, in CROPPED frame coords (image starts at
    # original (20, 200), so pill at original (80, 278) → cropped (60, 78)).
    pill_xy = (60, 78, 140, 24)
    px, py, pw, ph = pill_xy

    if n == 1:
        caption(d, 160, 24, "1. Find the commit selector pill on the Files toolbar.", f_cap)
        d.rectangle([px - 3, py - 3, px + pw + 3, py + ph + 3], outline=RED, width=4)
    elif n == 2:
        caption(d, 160, 24, "2. Click the selector — a dropdown lists every commit.", f_cap)
        d.rectangle([px - 3, py - 3, px + pw + 3, py + ph + 3], outline=RED, width=4)
        draw_dropdown(d, px, py + ph + 8, [
            ("All commits", "show everything (all 219 files)"),
            ("Scaffold v1.1 of RMC TotalRisk Users Guide", "commit 85e53a4 — 219 files of boilerplate"),
            ("Update preface, installation, GUI, and acronyms for v1.1", "commit 7ae2022 — the content edits"),
        ])
    elif n == 3:
        caption(d, 160, 24, "3. Pick the content commit to skip past the scaffolding.", f_cap)
        d.rectangle([px - 3, py - 3, px + pw + 3, py + ph + 3], outline=RED, width=4)
        draw_dropdown(d, px, py + ph + 8, [
            ("All commits", "show everything (all 219 files)"),
            ("Scaffold v1.1 of RMC TotalRisk Users Guide", "commit 85e53a4 — 219 files of boilerplate"),
            ("Update preface, installation, GUI, and acronyms for v1.1", "commit 7ae2022 — the content edits"),
        ], highlighted=2)
    elif n == 4:
        caption(d, 160, 24, "4. The file tree narrows to just the 4 files the content commit touches.", f_cap)
        # File tree in cropped frame: original (33, 384, 295, 160) → cropped (13, 184, 295, 160)
        d.rectangle([13, 184, 308, 344], outline=RED, width=4)

    return Image.alpha_composite(img.convert("RGBA"), overlay).convert("RGB")


def main():
    OUT_GIF.parent.mkdir(parents=True, exist_ok=True)
    OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)

    base = Image.open(SRC).convert("RGBA").crop(CROP)
    frames = [render_frame(base, i) for i in range(1, 5)]
    frames[0].save(OUT_GIF, save_all=True, append_images=frames[1:], duration=2400, loop=0, optimize=True)
    print(f"Wrote {OUT_GIF} ({base.size[0]}x{base.size[1]}, {len(frames)} frames)")

    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    prs.slide_width = Inches(base.size[0] / 96)
    prs.slide_height = Inches(base.size[1] / 96)
    blank = prs.slide_layouts[6]
    for frm in frames:
        slide = prs.slides.add_slide(blank)
        buf = io.BytesIO()
        frm.save(buf, format="PNG")
        buf.seek(0)
        slide.shapes.add_picture(buf, 0, 0, prs.slide_width, prs.slide_height)
    prs.save(str(OUT_PPTX))
    print(f"Wrote {OUT_PPTX} ({len(frames)} slides)")


if __name__ == "__main__":
    main()
