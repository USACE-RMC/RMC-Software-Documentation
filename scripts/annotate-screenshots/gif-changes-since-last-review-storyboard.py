"""
Storyboard GIF: viewing only what's new since your last review
(Section 7 — the backcheck round).

After the reviewer submits a Comment-style review and the author pushes
new commits, the Files-changed toolbar gains a "Changes since your last
review" option in the commit-range dropdown. This storyboard walks
through that flow.

Frames:
  1. After your first review, the author has pushed more commits
  2. Click the commit selector to open the dropdown
  3. The "Changes since your last review" option is at the top — pick it
  4. The diff narrows to only what's new

The state where this UI appears requires both a prior review and new
commits, which the sandbox PR doesn't naturally produce. The storyboard
mocks the dropdown overlay; the user re-records with ScreenToGif against
a real backcheck PR when one shows up.
"""

import io
from pathlib import Path
from PIL import Image, ImageDraw, ImageFont

REPO_ROOT = Path(__file__).resolve().parents[2]
SRC = REPO_ROOT / "planning" / "assets" / "captures" / "pr-files-changed.png"
OUT_GIF = REPO_ROOT / "static" / "figures" / "dev" / "reviewer-workflow" / "gif-changes-since-last-review-storyboard.gif"
OUT_PPTX = REPO_ROOT / "planning" / "assets" / "figure-pptx" / "gif-changes-since-last-review-storyboard.pptx"

CROP = (20, 200, 1340, 700)

ACCENT = (74, 124, 155, 255)
WHITE = (255, 255, 255, 255)
RED = (200, 30, 30, 255)
LABEL_BG = (255, 248, 220, 250)
BLUE_HOVER = (221, 244, 255, 255)


def load_font(size):
    for p in ("C:/Windows/Fonts/segoeuib.ttf", "C:/Windows/Fonts/arialbd.ttf"):
        if Path(p).exists():
            return ImageFont.truetype(p, size)
    return ImageFont.load_default()


def caption(draw, x, y, text, font):
    bb = draw.textbbox((0, 0), text, font=font)
    tw, th = bb[2] - bb[0], bb[3] - bb[1]
    pad = 12
    rect = [x, y, x + tw + pad * 2, y + th + pad * 2]
    draw.rounded_rectangle([rect[0] - 2, rect[1] - 2, rect[2] + 2, rect[3] + 2], radius=12, fill=WHITE)
    draw.rounded_rectangle(rect, radius=10, fill=LABEL_BG, outline=(50, 90, 115, 255), width=2)
    draw.text((x + pad - bb[0], y + pad - bb[1]), text, font=font, fill=(40, 40, 40, 255))


def draw_dropdown(draw, x, y, items, highlighted=None):
    item_h = 44
    w = 560
    h = item_h * len(items) + 16
    draw.rounded_rectangle([x, y, x + w, y + h], radius=8, fill=(255, 255, 255, 255), outline=(208, 215, 222, 255), width=1)
    f_title = load_font(13)
    f_sub = load_font(11)
    for i, (title, subtitle) in enumerate(items):
        iy = y + 8 + i * item_h
        if i == highlighted:
            draw.rectangle([x + 4, iy, x + w - 4, iy + item_h], fill=BLUE_HOVER)
        draw.text((x + 14, iy + 4), title, font=f_title, fill=(36, 41, 47, 255))
        draw.text((x + 14, iy + 22), subtitle, font=f_sub, fill=(101, 109, 118, 255))


def render_frame(base, n):
    img = base.copy()
    overlay = Image.new("RGBA", img.size, (0, 0, 0, 0))
    d = ImageDraw.Draw(overlay)
    f_cap = load_font(18)
    f_step = load_font(28)

    # Step indicator
    d.rounded_rectangle([20, 18, 130, 56], radius=8, fill=ACCENT)
    d.text((35, 24), f"Step {n}", font=f_step, fill=WHITE)

    pill_xy = (60, 78, 140, 24)  # commit selector in cropped frame
    px, py, pw, ph = pill_xy

    dropdown_items = [
        ("Changes since your last review", "only the new commits the author just pushed"),
        ("All commits", "everything in the PR"),
        ("Scaffold v1.1 of RMC TotalRisk Users Guide", "commit 85e53a4 — boilerplate"),
        ("Update preface, installation, GUI, and acronyms for v1.1", "commit 7ae2022 — the content edits"),
    ]

    if n == 1:
        caption(d, 160, 24, "1. You've already reviewed. The author has pushed new commits. Come back to the PR.", f_cap)
    elif n == 2:
        caption(d, 160, 24, "2. On Files changed, click the commit selector pill.", f_cap)
        d.rectangle([px - 3, py - 3, px + pw + 3, py + ph + 3], outline=RED, width=4)
    elif n == 3:
        caption(d, 160, 24, "3. Pick \"Changes since your last review\" at the top of the dropdown.", f_cap)
        d.rectangle([px - 3, py - 3, px + pw + 3, py + ph + 3], outline=RED, width=4)
        draw_dropdown(d, px, py + ph + 8, dropdown_items, highlighted=0)
    elif n == 4:
        caption(d, 160, 24, "4. The diff narrows to only the new changes — backcheck just those.", f_cap)
        # Highlight the file tree + diff area
        d.rectangle([13, 184, 308, 344], outline=RED, width=4)
        d.rectangle([336, 380, 1316, 470], outline=RED, width=4)

    return Image.alpha_composite(img.convert("RGBA"), overlay).convert("RGB")


def main():
    OUT_GIF.parent.mkdir(parents=True, exist_ok=True)
    OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)

    base = Image.open(SRC).convert("RGBA").crop(CROP)
    frames = [render_frame(base, i) for i in range(1, 5)]
    frames[0].save(OUT_GIF, save_all=True, append_images=frames[1:], duration=2400, loop=0, optimize=True)
    print(f"Wrote {OUT_GIF} ({base.size[0]}x{base.size[1]}, {len(frames)} frames)")

    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    prs.slide_width = Inches(base.size[0] / 96)
    prs.slide_height = Inches(base.size[1] / 96)
    blank = prs.slide_layouts[6]
    for frm in frames:
        slide = prs.slides.add_slide(blank)
        buf = io.BytesIO()
        frm.save(buf, format="PNG")
        buf.seek(0)
        slide.shapes.add_picture(buf, 0, 0, prs.slide_width, prs.slide_height)
    prs.save(str(OUT_PPTX))
    print(f"Wrote {OUT_PPTX} ({len(frames)} slides)")


if __name__ == "__main__":
    main()
