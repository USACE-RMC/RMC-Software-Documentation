"""
Storyboard GIF for the comment-composer interaction (Section 5 of the
reviewer-workflow chapter).

Why a storyboard, not a real-capture GIF: GitHub's new diff grid mounts
the per-line "Add a comment" button via React event handlers that only
fire on real cursor motion — Playwright's synthetic hover events don't
trigger them. Until that's resolved, this script produces a stepped GIF
from the static Files-changed capture by overlaying per-frame callouts.
The user re-records this with ScreenToGif against the live PR.

Frames:
  1. Hover a line in the diff
  2. Blue + button appears
  3. Click + → composer opens
  4. Optional: click ± for a suggested change
  5. Click Start a review

Outputs:
  static/figures/dev/reviewer-workflow/gif-composer-storyboard.gif
  planning/assets/figure-pptx/gif-composer-storyboard.pptx
    (the pptx contains five slides — one per frame — so the user can
    customize the storyboard before re-recording)
"""

import io
from pathlib import Path
from PIL import Image, ImageDraw, ImageFont

REPO_ROOT = Path(__file__).resolve().parents[2]
SRC = REPO_ROOT / "planning" / "assets" / "captures" / "pr-files-changed.png"
OUT_GIF = REPO_ROOT / "static" / "figures" / "dev" / "reviewer-workflow" / "gif-composer-storyboard.gif"
OUT_PPTX = REPO_ROOT / "planning" / "assets" / "figure-pptx" / "gif-composer-storyboard.pptx"

# Crop to a region that contains the file header + a few lines of the
# 01-preface diff. The captured page is enormous (61k px tall) so we
# zoom in on the relevant section near the top.
CROP = (20, 200, 1340, 950)

ACCENT = (74, 124, 155, 255)
ACCENT_DARK = (50, 90, 115, 255)
RED = (200, 30, 30, 255)
LABEL_BG = (255, 248, 220, 250)
WHITE = (255, 255, 255, 255)


def load_font(size):
    for p in ("C:/Windows/Fonts/segoeuib.ttf", "C:/Windows/Fonts/arialbd.ttf"):
        if Path(p).exists():
            return ImageFont.truetype(p, size)
    return ImageFont.load_default()


def caption(draw, x, y, text, font, max_width=620):
    """A soft-yellow rounded caption box anchored at (x, y) top-left."""
    bb = draw.textbbox((0, 0), text, font=font)
    tw, th = bb[2] - bb[0], bb[3] - bb[1]
    pad = 12
    rect = [x, y, x + tw + pad * 2, y + th + pad * 2]
    draw.rounded_rectangle([rect[0] - 2, rect[1] - 2, rect[2] + 2, rect[3] + 2], radius=12, fill=WHITE)
    draw.rounded_rectangle(rect, radius=10, fill=LABEL_BG, outline=ACCENT_DARK, width=2)
    draw.text((x + pad - bb[0], y + pad - bb[1]), text, font=font, fill=(40, 40, 40, 255))


def draw_arrow(draw, from_xy, to_xy, color=RED, width=4):
    fx, fy = from_xy
    tx, ty = to_xy
    draw.line([fx, fy, tx, ty], fill=color, width=width)
    # Simple arrowhead: small triangle at the target end
    import math
    angle = math.atan2(ty - fy, tx - fx)
    ahead = 14
    ax = tx - ahead * math.cos(angle)
    ay = ty - ahead * math.sin(angle)
    p1 = (ax + ahead * 0.5 * math.cos(angle + math.pi / 2),
          ay + ahead * 0.5 * math.sin(angle + math.pi / 2))
    p2 = (ax - ahead * 0.5 * math.cos(angle + math.pi / 2),
          ay - ahead * 0.5 * math.sin(angle + math.pi / 2))
    draw.polygon([(tx, ty), p1, p2], fill=color)


def draw_blue_plus(draw, x, y, size=22):
    """Draw a faux 'add comment' blue plus badge, mimicking GitHub's hover icon."""
    draw.rounded_rectangle([x - size // 2, y - size // 2, x + size // 2, y + size // 2], radius=6, fill=(31, 111, 235, 255))
    # plus
    cx, cy = x, y
    arm = size // 2 - 6
    draw.line([cx - arm, cy, cx + arm, cy], fill=WHITE, width=3)
    draw.line([cx, cy - arm, cx, cy + arm], fill=WHITE, width=3)


def draw_composer_mock(draw, x, y, w, h, with_suggestion=False, with_typed_text=False):
    """Approximate the inline-comment composer."""
    # Outer card
    draw.rounded_rectangle([x, y, x + w, y + h], radius=8, fill=(255, 255, 255, 255), outline=(208, 215, 222, 255), width=1)
    # Toolbar tabs (Write / Preview)
    draw.rectangle([x, y, x + w, y + 32], fill=(246, 248, 250, 255))
    fbold = load_font(13)
    draw.text((x + 14, y + 9), "Write", font=fbold, fill=(36, 41, 47, 255))
    draw.text((x + 70, y + 9), "Preview", font=fbold, fill=(101, 109, 118, 255))
    # Markdown toolbar
    fsmall = load_font(12)
    icons = "H  B  I  </>  •••  @  ↵  ±"
    draw.text((x + 14, y + 42), icons, font=fsmall, fill=(101, 109, 118, 255))
    # Text area
    draw.rectangle([x + 8, y + 70, x + w - 8, y + h - 60], fill=(255, 255, 255, 255), outline=(208, 215, 222, 255), width=1)
    if with_typed_text:
        draw.text((x + 16, y + 80), 'Consider clarifying which "refinements" are meant — see line 24.', font=fsmall, fill=(36, 41, 47, 255))
    if with_suggestion:
        # Suggestion code block
        sx, sy = x + 16, y + 100
        sw, sh = w - 32, 70
        draw.rectangle([sx, sy, sx + sw, sy + sh], fill=(246, 248, 250, 255), outline=(208, 215, 222, 255), width=1)
        draw.text((sx + 8, sy + 6), "```suggestion", font=fsmall, fill=(101, 109, 118, 255))
        draw.text((sx + 8, sy + 24), 'Version 1.1 of this Users Guide documents the v1.1 release of', font=fsmall, fill=(36, 41, 47, 255))
        draw.text((sx + 8, sy + 40), 'RMC-TotalRisk.', font=fsmall, fill=(36, 41, 47, 255))
        draw.text((sx + 8, sy + 54), "```", font=fsmall, fill=(101, 109, 118, 255))
    # Bottom-right buttons
    btn_y = y + h - 38
    draw.rounded_rectangle([x + w - 280, btn_y, x + w - 175, btn_y + 28], radius=4, fill=(246, 248, 250, 255), outline=(208, 215, 222, 255), width=1)
    draw.text((x + w - 258, btn_y + 7), "Cancel", font=fsmall, fill=(36, 41, 47, 255))
    draw.rounded_rectangle([x + w - 165, btn_y, x + w - 12, btn_y + 28], radius=4, fill=(46, 164, 79, 255))
    draw.text((x + w - 152, btn_y + 7), "Start a review", font=fsmall, fill=WHITE)


def render_frame(base: Image.Image, n: int) -> Image.Image:
    img = base.copy()
    overlay = Image.new("RGBA", img.size, (0, 0, 0, 0))
    d = ImageDraw.Draw(overlay)
    f_caption = load_font(18)
    f_step = load_font(28)

    # Step indicator (top-left)
    d.rounded_rectangle([20, 18, 130, 56], radius=8, fill=ACCENT)
    d.text((35, 24), f"Step {n}", font=f_step, fill=WHITE)

    # Line 23 in the diff sits at roughly (370, 778) in CROP coords —
    # this is the gutter cell for the first added preface line.
    line_xy = (370, 778)
    gutter_x, gutter_y = 350, 778

    if n == 1:
        caption(d, 160, 24, "Hover the line number gutter where you want to comment.", f_caption)
        draw_arrow(d, (550, 80), (gutter_x + 5, gutter_y - 4), color=RED)
    elif n == 2:
        caption(d, 160, 24, "A blue + button appears on the line. Click it.", f_caption)
        draw_blue_plus(d, gutter_x - 12, gutter_y + 4)
        draw_arrow(d, (550, 80), (gutter_x - 12, gutter_y + 4), color=RED)
    elif n == 3:
        caption(d, 160, 24, "The composer opens inline. Type your feedback.", f_caption)
        draw_composer_mock(d, 350, 820, 980, 220, with_typed_text=True)
    elif n == 4:
        caption(d, 160, 24, "Optional: click the ± icon to write a suggested change.", f_caption)
        draw_composer_mock(d, 350, 820, 980, 240, with_typed_text=True, with_suggestion=True)
        # Highlight the ± in the mock toolbar
        d.ellipse([558, 858, 590, 884], outline=RED, width=3)
    elif n == 5:
        caption(d, 160, 24, "Click \"Start a review\" to batch this with later comments.", f_caption)
        draw_composer_mock(d, 350, 820, 980, 240, with_typed_text=True, with_suggestion=True)
        # Highlight the Start a review button
        d.rectangle([1162, 1018, 1318, 1050], outline=RED, width=4)

    composed = Image.alpha_composite(img.convert("RGBA"), overlay).convert("RGB")
    return composed


def main():
    OUT_GIF.parent.mkdir(parents=True, exist_ok=True)
    OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)

    base = Image.open(SRC).convert("RGBA").crop(CROP)
    frames = [render_frame(base, i) for i in range(1, 6)]
    frames[0].save(
        OUT_GIF,
        save_all=True,
        append_images=frames[1:],
        duration=2200,
        loop=0,
        optimize=True,
    )
    print(f"Wrote {OUT_GIF} ({base.size[0]}x{base.size[1]}, 5 frames)")

    # PPTX: one slide per frame, with the rendered frame as background.
    # Annotations are baked into the frame image for the storyboard
    # (since the storyboard is a guide for re-recording, not a live
    # figure the user iterates on shape positions for).
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    prs.slide_width = Inches(base.size[0] / 96)
    prs.slide_height = Inches(base.size[1] / 96)
    blank = prs.slide_layouts[6]
    for i, frm in enumerate(frames, start=1):
        slide = prs.slides.add_slide(blank)
        buf = io.BytesIO()
        frm.save(buf, format="PNG")
        buf.seek(0)
        slide.shapes.add_picture(buf, 0, 0, prs.slide_width, prs.slide_height)
    prs.save(str(OUT_PPTX))
    print(f"Wrote {OUT_PPTX} ({len(frames)} slides)")


if __name__ == "__main__":
    main()
